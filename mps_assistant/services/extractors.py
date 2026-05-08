from __future__ import annotations

import csv
import hashlib
import json
import re
import subprocess
from pathlib import Path
from typing import Any, List, Optional, Sequence, Tuple

from bs4 import BeautifulSoup
from docx import Document
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation
import xlrd

from ..schemas import ExtractedDocument, ExtractedSection


def build_extracted_document(
    *,
    source_key: str,
    origin: str,
    source_format: str,
    downloaded_at: str,
    sections: Sequence[ExtractedSection],
    checksum_content: Any,
    url: Optional[str] = None,
    local_path: Optional[str] = None,
    page_title: Optional[str] = None,
    document_title: Optional[str] = None,
    file_name: Optional[str] = None,
    content_type: Optional[str] = None,
) -> ExtractedDocument:
    checksum_bytes = _checksum_bytes_from_content(checksum_content)
    return ExtractedDocument(
        source_key=source_key,
        origin=origin,
        source_format=source_format,
        url=url,
        local_path=local_path,
        page_title=page_title,
        document_title=document_title,
        file_name=file_name,
        content_type=content_type,
        downloaded_at=downloaded_at,
        checksum=_sha256_bytes(checksum_bytes),
        sections=list(sections),
    )


def extract_json_document(
    *,
    source_key: str,
    origin: str,
    downloaded_at: str,
    payload: Any,
    url: Optional[str] = None,
    page_title: Optional[str] = None,
    document_title: Optional[str] = None,
    file_name: Optional[str] = None,
    content_type: str = "application/json",
    source_format: str = "json",
) -> ExtractedDocument:
    title = document_title or page_title or file_name or source_key
    sections = _json_sections(payload, title)
    return build_extracted_document(
        source_key=source_key,
        origin=origin,
        source_format=source_format,
        downloaded_at=downloaded_at,
        sections=sections,
        checksum_content=payload,
        url=url,
        local_path=None,
        page_title=page_title,
        document_title=title,
        file_name=file_name,
        content_type=content_type,
    )


def extract_html_document(
    url: str,
    html: str,
    downloaded_at: str,
    *,
    page_title_override: Optional[str] = None,
    document_title_override: Optional[str] = None,
    source_format: str = "html",
) -> Tuple[ExtractedDocument, List[str]]:
    soup = BeautifulSoup(html, "html.parser")
    heading_title_hint = _first_text(soup.select_one(".titlePage, .title-page")) or _meaningful_heading_from_soup(soup)
    for tag in soup(["script", "style", "noscript", "svg", "header", "footer", "nav"]):
        tag.decompose()

    _remove_cookie_banners(soup)
    page_title = page_title_override or _get_text(soup.title)
    main = _select_content_root(soup)
    document_title = (
        document_title_override
        or
        _first_text(main.find(["h1", "h2", "h3"]))
        or heading_title_hint
        or _first_text(soup.find(["h1", "h2", "h3"]))
        or page_title
    )
    document_title = _normalize_document_title(document_title)
    links = [anchor.get("href", "").strip() for anchor in soup.find_all("a", href=True)]
    sections: List[ExtractedSection]
    if _has_meaningful_form(main):
        non_form_main = BeautifulSoup(str(main), "html.parser")
        for form in non_form_main.find_all("form"):
            form.decompose()
        sections = _extract_html_sections(non_form_main, document_title)
        sections.extend(_extract_form_sections(main, document_title))
        sections = _dedupe_sections(sections)
    else:
        sections = _extract_html_sections(main, document_title)

    document = ExtractedDocument(
        source_key=url,
        origin="website",
        source_format=source_format,
        url=url,
        local_path=None,
        page_title=page_title,
        document_title=document_title,
        file_name=None,
        content_type="text/html",
        downloaded_at=downloaded_at,
        checksum=_sha256_bytes(html.encode("utf-8", errors="ignore")),
        sections=sections,
    )
    return document, links


def extract_file_document(
    source_key: str,
    origin: str,
    path: Path,
    downloaded_at: str,
    url: Optional[str] = None,
    content_type: Optional[str] = None,
) -> ExtractedDocument:
    suffix = path.suffix.lower()
    file_name = path.name

    if suffix == ".pdf":
        sections, title = _extract_pdf(path)
        source_format = "pdf"
    elif suffix == ".docx":
        sections, title = _extract_docx(path)
        source_format = "docx"
    elif suffix == ".doc":
        sections, title = _extract_textutil_document(path)
        source_format = "doc"
    elif suffix == ".xlsx":
        sections, title = _extract_xlsx(path)
        source_format = "xlsx"
    elif suffix == ".xls":
        sections, title = _extract_xls(path)
        source_format = "xls"
    elif suffix in {".ppt", ".pptx"}:
        sections, title = _extract_pptx(path)
        source_format = "pptx"
    elif suffix == ".csv":
        sections, title = _extract_csv(path)
        source_format = "csv"
    elif suffix == ".rtf":
        sections, title = _extract_textutil_document(path)
        source_format = "rtf"
    else:
        sections, title = _extract_text_file(path)
        source_format = suffix.lstrip(".") or "text"

    return ExtractedDocument(
        source_key=source_key,
        origin=origin,
        source_format=source_format,
        url=url,
        local_path=str(path),
        page_title=title,
        document_title=title,
        file_name=file_name,
        content_type=content_type,
        downloaded_at=downloaded_at,
        checksum=_sha256_bytes(path.read_bytes()),
        sections=sections,
    )


def _extract_html_sections(root: BeautifulSoup, fallback_heading: Optional[str]) -> List[ExtractedSection]:
    sections: List[ExtractedSection] = []
    current_heading = fallback_heading
    buffer: List[str] = []

    for element in root.find_all(["h1", "h2", "h3", "h4", "h5", "h6", "p", "li", "table"]):
        if element.name and element.name.startswith("h"):
            _flush_html_section(sections, current_heading, buffer)
            current_heading = _get_text(element) or current_heading
            buffer = []
            continue

        if element.name == "table":
            table_text = _table_to_text(element)
            if table_text:
                buffer.append(table_text)
            continue

        text = _get_text(element)
        if text:
            buffer.append(text)

    _flush_html_section(sections, current_heading, buffer)

    if not sections:
        text = _get_text(root)
        if text:
            sections.append(ExtractedSection(heading=fallback_heading, text=text))

    return sections


def _extract_form_sections(root: BeautifulSoup, fallback_heading: Optional[str]) -> List[ExtractedSection]:
    sections: List[ExtractedSection] = []
    forms = root.find_all("form")
    for form_index, form in enumerate(forms, start=1):
        containers = form.select(".card, fieldset")
        if not containers:
            containers = [form]

        for container_index, container in enumerate(containers, start=1):
            heading = (
                _first_text(container.find(class_="card-header"))
                or _first_text(container.find(["legend", "h1", "h2", "h3", "h4"]))
                or fallback_heading
                or f"Form section {form_index}.{container_index}"
            )
            text = _container_text(container)
            if text and not _is_noise_section(heading, text):
                sections.append(ExtractedSection(heading=heading, text=text))
    return sections


def _has_meaningful_form(root: BeautifulSoup) -> bool:
    for form in root.find_all("form"):
        labels = form.find_all(["label", "legend"])
        controls = []
        for control in form.find_all(["input", "select", "textarea"]):
            if control.name == "input":
                control_type = (control.get("type") or "text").lower()
                if control_type in {"hidden", "submit", "button", "image"}:
                    continue
            controls.append(control)
        if labels and controls:
            return True
        if len(controls) >= 4:
            return True
    return False


def _select_content_root(soup: BeautifulSoup):
    candidates = [
        soup.find("main"),
        soup.find(id="main"),
        soup.find("article"),
        soup.body,
        soup,
    ]
    for candidate in candidates:
        if candidate is None:
            continue
        if _candidate_has_content(candidate):
            return candidate
    return soup.body or soup


def _candidate_has_content(candidate) -> bool:
    if candidate.find(["h1", "h2", "h3", "p", "li", "table", "form"]):
        return True
    return len(_get_text(candidate)) >= 120


def _meaningful_heading_from_soup(soup: BeautifulSoup) -> Optional[str]:
    ignored = {
        "manage consent preferences",
        "privacy preference center",
        "cookies button",
        "south africa",
    }
    for heading in soup.find_all(["h1", "h2", "h3"]):
        text = _get_text(heading)
        if not text:
            continue
        if text.strip().lower() in ignored:
            continue
        return text
    return None


def _normalize_document_title(title: Optional[str]) -> str:
    value = (title or "").strip()
    if not value:
        return ""
    if value.lower().startswith("to successfully complete this application"):
        return "Membership application"
    return value


def _flush_html_section(sections: List[ExtractedSection], heading: Optional[str], buffer: List[str]) -> None:
    if not buffer:
        return
    text = "\n\n".join(part for part in buffer if part)
    if text.strip() and not _is_noise_section(heading, text):
        sections.append(ExtractedSection(heading=heading, text=text.strip()))


def _extract_pdf(path: Path) -> Tuple[List[ExtractedSection], str]:
    reader = PdfReader(str(path))
    title = getattr(reader.metadata, "title", None) or path.stem
    sections: List[ExtractedSection] = []
    for index, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        text = text.strip()
        if text:
            sections.append(ExtractedSection(heading=title, text=text, page_number=index))
    return sections, title


def _extract_docx(path: Path) -> Tuple[List[ExtractedSection], str]:
    doc = Document(str(path))
    title = doc.core_properties.title or path.stem
    sections: List[ExtractedSection] = []
    current_heading = title
    buffer: List[str] = []
    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style_name = getattr(paragraph.style, "name", "") or ""
        if style_name.lower().startswith("heading"):
            if buffer:
                sections.append(ExtractedSection(heading=current_heading, text="\n\n".join(buffer)))
                buffer = []
            current_heading = text
        else:
            buffer.append(text)
    if buffer:
        sections.append(ExtractedSection(heading=current_heading, text="\n\n".join(buffer)))
    return sections, title


def _extract_xlsx(path: Path) -> Tuple[List[ExtractedSection], str]:
    workbook = load_workbook(filename=str(path), read_only=True, data_only=True)
    title = path.stem
    sections: List[ExtractedSection] = []
    for sheet in workbook.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            values = [str(value).strip() for value in row if value not in (None, "")]
            if values:
                rows.append(" | ".join(values))
        if rows:
            sections.append(ExtractedSection(heading=sheet.title, text="\n".join(rows)))
    return sections, title


def _extract_xls(path: Path) -> Tuple[List[ExtractedSection], str]:
    workbook = xlrd.open_workbook(str(path))
    title = path.stem
    sections: List[ExtractedSection] = []
    for sheet in workbook.sheets():
        rows = []
        for row_index in range(sheet.nrows):
            values = [str(value).strip() for value in sheet.row_values(row_index) if str(value).strip()]
            if values:
                rows.append(" | ".join(values))
        if rows:
            sections.append(ExtractedSection(heading=sheet.name, text="\n".join(rows)))
    return sections, title


def _extract_pptx(path: Path) -> Tuple[List[ExtractedSection], str]:
    presentation = Presentation(str(path))
    title = path.stem
    sections: List[ExtractedSection] = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text.strip())
        text = "\n\n".join(part for part in texts if part)
        if text:
            sections.append(ExtractedSection(heading=f"Slide {index}", text=text, page_number=index))
    return sections, title


def _extract_csv(path: Path) -> Tuple[List[ExtractedSection], str]:
    title = path.stem
    rows = []
    with path.open("r", encoding="utf-8", errors="ignore", newline="") as handle:
        reader = csv.reader(handle)
        for row in reader:
            cleaned = [value.strip() for value in row if value and value.strip()]
            if cleaned:
                rows.append(" | ".join(cleaned))
    sections = [ExtractedSection(heading=title, text="\n".join(rows))] if rows else []
    return sections, title


def _extract_text_file(path: Path) -> Tuple[List[ExtractedSection], str]:
    title = path.stem
    text = path.read_text(encoding="utf-8", errors="ignore").strip()
    sections = [ExtractedSection(heading=title, text=text)] if text else []
    return sections, title


def _extract_textutil_document(path: Path) -> Tuple[List[ExtractedSection], str]:
    title = path.stem
    try:
        result = subprocess.run(
            ["textutil", "-convert", "txt", "-stdout", str(path)],
            capture_output=True,
            check=True,
        )
        text = result.stdout.decode("utf-8", errors="ignore").strip()
    except Exception:
        text = path.read_text(encoding="utf-8", errors="ignore").strip()
    sections = [ExtractedSection(heading=title, text=text)] if text else []
    return sections, title


def _table_to_text(table) -> str:
    rows = []
    for row in table.find_all("tr"):
        cells = [cell.get_text(" ", strip=True) for cell in row.find_all(["td", "th"])]
        cells = [cell for cell in cells if cell]
        if cells:
            rows.append(" | ".join(cells))
    return "\n".join(rows)


def _container_text(container) -> str:
    ignored_lines = {
        "accept",
        "reject all",
        "let me choose",
        "cookies button",
        "back",
        "submit",
    }
    lines = []
    seen = set()
    for fragment in container.stripped_strings:
        line = " ".join(fragment.split())
        if not line:
            continue
        if line.lower() in ignored_lines:
            continue
        if line.lower().startswith("by clicking “accept”"):
            continue
        if line not in seen:
            seen.add(line)
            lines.append(line)
    return "\n".join(lines).strip()


def _dedupe_sections(sections: List[ExtractedSection]) -> List[ExtractedSection]:
    deduped: List[ExtractedSection] = []
    seen = set()
    for section in sections:
        key = (section.heading or "", " ".join(section.text.split()))
        if key in seen:
            continue
        seen.add(key)
        deduped.append(section)
    return deduped


def _remove_cookie_banners(soup: BeautifulSoup) -> None:
    for element in soup.find_all(True):
        if getattr(element, "attrs", None) is None:
            continue
        classes = " ".join(element.get("class", [])).lower()
        attrs = " ".join(f"{key}={value}" for key, value in element.attrs.items()).lower()
        text = " ".join(element.get_text(" ", strip=True).split()).lower()
        if "cookie" in classes or "cookie" in attrs:
            element.decompose()
            continue
        if text.startswith("by clicking “accept”") or text.startswith('by clicking "accept"'):
            element.decompose()


def _is_noise_section(heading: Optional[str], text: str) -> bool:
    heading_text = (heading or "").strip().lower()
    text_value = " ".join(text.strip().lower().split())
    if heading_text in {"share this page", "related tags"}:
        return True
    if text_value in {
        "print email facebook twitter google+ linked in",
        "facebook twitter google+ linked in",
    }:
        return True
    if "by clicking “accept”" in text_value or 'by clicking "accept"' in text_value:
        return True
    return False


def _first_text(element) -> Optional[str]:
    if element is None:
        return None
    text = element.get_text(" ", strip=True)
    return text or None


def _get_text(element) -> str:
    if element is None:
        return ""
    return " ".join(element.get_text(" ", strip=True).split())


def _sha256_bytes(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def _checksum_bytes_from_content(content: Any) -> bytes:
    if isinstance(content, bytes):
        return content
    if isinstance(content, str):
        return content.encode("utf-8", errors="ignore")
    return json.dumps(content, ensure_ascii=True, sort_keys=True).encode("utf-8")


def _json_sections(payload: Any, title: str) -> List[ExtractedSection]:
    if isinstance(payload, dict):
        sections = []
        scalar_lines = []
        for key, value in payload.items():
            if _is_scalar_json_value(value):
                scalar_lines.append(f"{_friendly_key(key)}: {_scalar_value_text(value)}")
                continue
            section_text = _structured_json_text(value)
            if section_text:
                sections.append(ExtractedSection(heading=_friendly_key(key), text=section_text))
        if scalar_lines:
            sections.insert(0, ExtractedSection(heading=title, text="\n".join(scalar_lines)))
        return sections

    section_text = _structured_json_text(payload)
    return [ExtractedSection(heading=title, text=section_text)] if section_text else []


def _structured_json_text(value: Any, indent: int = 0) -> str:
    prefix = "  " * indent
    if isinstance(value, dict):
        lines = []
        for key, item in value.items():
            label = _friendly_key(key)
            if _is_scalar_json_value(item):
                lines.append(f"{prefix}{label}: {_scalar_value_text(item)}")
            else:
                nested = _structured_json_text(item, indent + 1)
                if nested:
                    lines.append(f"{prefix}{label}:")
                    lines.append(nested)
        return "\n".join(lines).strip()

    if isinstance(value, list):
        lines = []
        for item in value:
            if _is_scalar_json_value(item):
                lines.append(f"{prefix}- {_scalar_value_text(item)}")
            else:
                nested = _structured_json_text(item, indent + 1)
                if nested:
                    lines.append(f"{prefix}-")
                    lines.append(nested)
        return "\n".join(lines).strip()

    return f"{prefix}{_scalar_value_text(value)}".strip()


def _friendly_key(value: str) -> str:
    spaced = re.sub(r"(?<!^)(?=[A-Z])", " ", value or "")
    spaced = spaced.replace("_", " ").replace("-", " ")
    return " ".join(spaced.split()).strip().capitalize() or value


def _is_scalar_json_value(value: Any) -> bool:
    return value is None or isinstance(value, (str, int, float, bool))


def _scalar_value_text(value: Any) -> str:
    if value is None:
        return "None"
    if isinstance(value, bool):
        return "Yes" if value else "No"
    return str(value)
